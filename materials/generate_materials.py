"""根据 Markdown 源稿生成教师 Word 指导手册和课堂 PPT。"""

from pathlib import Path


def _add_code_block(document, lines):
    paragraph = document.add_paragraph()
    run = paragraph.add_run('\n'.join(lines))
    run.font.name = 'Consolas'


def _add_markdown_table(document, rows):
    parsed = [
        [cell.strip() for cell in row.strip().strip('|').split('|')]
        for row in rows
        if not set(row.replace('|', '').replace(' ', '').strip()) <= {'-', ':'}
    ]
    if not parsed:
        return
    table = document.add_table(rows=len(parsed), cols=len(parsed[0]))
    table.style = 'Table Grid'
    for row_index, row in enumerate(parsed):
        for column_index, value in enumerate(row):
            table.cell(row_index, column_index).text = value


def generate_word():
    from docx import Document

    source = Path('materials/teacher_lab_guide.md')
    document = Document()
    in_code = False
    code_lines = []
    table_lines = []

    def flush_table():
        nonlocal table_lines
        if table_lines:
            _add_markdown_table(document, table_lines)
            table_lines = []

    for line in source.read_text(encoding='utf-8').splitlines():
        if line.startswith('```'):
            flush_table()
            if in_code:
                _add_code_block(document, code_lines)
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_lines.append(line)
            continue
        if line.strip().startswith('|') and line.strip().endswith('|'):
            table_lines.append(line)
            continue

        flush_table()
        if line.startswith('# '):
            document.add_heading(line[2:], level=0)
        elif line.startswith('## '):
            document.add_heading(line[3:], level=1)
        elif line.startswith('### '):
            document.add_heading(line[4:], level=2)
        elif line.startswith('#### '):
            document.add_heading(line[5:], level=3)
        elif line.startswith('- '):
            document.add_paragraph(line[2:], style='List Bullet')
        elif len(line) > 3 and line[0].isdigit() and '. ' in line[:5]:
            document.add_paragraph(line.split('. ', 1)[1], style='List Number')
        elif line.strip() == '---':
            document.add_paragraph('')
        elif line.strip():
            document.add_paragraph(line)
    flush_table()
    if code_lines:
        _add_code_block(document, code_lines)
    target = Path('materials/teacher_lab_guide.docx')
    try:
        document.save(target)
    except PermissionError:
        fallback = Path('materials/teacher_lab_guide_student_v2.docx')
        document.save(fallback)
        print(f'{target} 正在被 Word/WPS 占用，已生成备用文件：{fallback}')


def generate_ppt():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    colors = {
        'navy': RGBColor(14, 33, 66),
        'blue': RGBColor(37, 99, 235),
        'cyan': RGBColor(14, 165, 233),
        'green': RGBColor(16, 185, 129),
        'amber': RGBColor(245, 158, 11),
        'red': RGBColor(239, 68, 68),
        'slate': RGBColor(71, 85, 105),
        'muted': RGBColor(100, 116, 139),
        'light': RGBColor(241, 245, 249),
        'white': RGBColor(255, 255, 255),
        'dark': RGBColor(15, 23, 42),
    }

    def set_text_frame(shape, text, font_size=20, color=None, bold=False, align=None):
        shape.text = text
        paragraph = shape.text_frame.paragraphs[0]
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = 'Microsoft YaHei'
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color

    def add_textbox(slide, left, top, width, height, text, font_size=20, color=None, bold=False, align=None):
        shape = slide.shapes.add_textbox(left, top, width, height)
        shape.text_frame.margin_left = Inches(0.05)
        shape.text_frame.margin_right = Inches(0.05)
        shape.text_frame.margin_top = Inches(0.02)
        shape.text_frame.margin_bottom = Inches(0.02)
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text_frame(shape, text, font_size=font_size, color=color, bold=bold, align=align)
        return shape

    def add_round_rect(slide, left, top, width, height, fill, line=None, radius=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = line or fill
        return shape

    def add_header(slide, title, slide_number, total):
        add_round_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.72), colors['navy'], radius=False)
        add_textbox(slide, Inches(0.45), Inches(0.12), Inches(10.4), Inches(0.42), title, 22, colors['white'], True)
        add_textbox(slide, Inches(11.4), Inches(0.18), Inches(1.35), Inches(0.32), f'{slide_number:02d}/{total:02d}', 12, colors['white'], False, PP_ALIGN.RIGHT)

    def add_footer(slide):
        add_round_rect(slide, Inches(0.45), Inches(7.08), Inches(12.45), Inches(0.02), colors['cyan'], radius=False)
        add_textbox(slide, Inches(0.45), Inches(7.12), Inches(6.0), Inches(0.25), '无线通信技术 · 信道编码与信道均衡实验', 9, colors['muted'])
        add_textbox(slide, Inches(8.5), Inches(7.12), Inches(4.4), Inches(0.25), '代码补全 · 仿真验证 · GitHub 自动评分', 9, colors['muted'], False, PP_ALIGN.RIGHT)

    def add_bullet_cards(slide, bullets, left=0.65, top=1.15, width=6.15, gap=0.12):
        card_height = min(0.73, max(0.52, (5.65 - gap * (len(bullets) - 1)) / max(len(bullets), 1)))
        accent_cycle = [colors['blue'], colors['cyan'], colors['green'], colors['amber'], colors['red'], colors['slate']]
        for index, bullet in enumerate(bullets[:6]):
            card_top = Inches(top + index * (card_height + gap))
            add_round_rect(slide, Inches(left), card_top, Inches(width), Inches(card_height), colors['white'], RGBColor(226, 232, 240))
            add_round_rect(slide, Inches(left), card_top, Inches(0.11), Inches(card_height), accent_cycle[index % len(accent_cycle)], radius=False)
            font_size = 17 if len(bullet) <= 36 else 14
            add_textbox(slide, Inches(left + 0.22), card_top + Inches(0.06), Inches(width - 0.35), Inches(card_height - 0.12), bullet, font_size, colors['dark'])

    def add_badge(slide, left, top, text, fill, width=1.45):
        add_round_rect(slide, Inches(left), Inches(top), Inches(width), Inches(0.42), fill)
        add_textbox(slide, Inches(left + 0.05), Inches(top + 0.06), Inches(width - 0.1), Inches(0.26), text, 11, colors['white'], True, PP_ALIGN.CENTER)

    def add_visual_panel(slide, title, bullets, slide_number):
        panel_left = Inches(7.1)
        panel_top = Inches(1.15)
        panel_width = Inches(5.55)
        panel_height = Inches(5.65)
        add_round_rect(slide, panel_left, panel_top, panel_width, panel_height, colors['light'], RGBColor(226, 232, 240))
        lower_title = title.lower()

        if slide_number == 2:
            add_textbox(slide, Inches(7.45), Inches(1.45), Inches(4.9), Inches(0.36), '课堂节奏', 18, colors['navy'], True, PP_ALIGN.CENTER)
            stages = [('0-5', '问题'), ('5-13', '编码'), ('13-23', '均衡'), ('23-30', '提交')]
            for index, (time_text, label) in enumerate(stages):
                y = 2.1 + index * 0.9
                add_round_rect(slide, Inches(7.55), Inches(y), Inches(1.15), Inches(0.45), colors['blue'])
                add_textbox(slide, Inches(7.6), Inches(y + 0.08), Inches(1.05), Inches(0.25), time_text, 12, colors['white'], True, PP_ALIGN.CENTER)
                add_textbox(slide, Inches(8.95), Inches(y + 0.04), Inches(2.9), Inches(0.34), label, 18, colors['dark'], True)
        elif '无线链路' in title:
            nodes = ['Bits', 'Encode', 'Channel', 'Equalize', 'Decode']
            for index, node in enumerate(nodes):
                x = 7.35 + index * 0.98
                add_round_rect(slide, Inches(x), Inches(2.6), Inches(0.82), Inches(0.56), colors['white'], RGBColor(203, 213, 225))
                add_textbox(slide, Inches(x + 0.03), Inches(2.75), Inches(0.76), Inches(0.18), node, 8, colors['navy'], True, PP_ALIGN.CENTER)
                if index < len(nodes) - 1:
                    add_textbox(slide, Inches(x + 0.78), Inches(2.72), Inches(0.28), Inches(0.22), '→', 18, colors['cyan'], True, PP_ALIGN.CENTER)
            add_badge(slide, 7.65, 4.15, '误码', colors['red'])
            add_badge(slide, 9.55, 4.15, 'ISI', colors['amber'])
            add_badge(slide, 11.05, 4.15, '补偿', colors['green'])
        elif 'hamming' in lower_title or '矩阵' in title or '校验' in title or '纠错' in title:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), 'Hamming(7,4) 核心链路', 17, colors['navy'], True, PP_ALIGN.CENTER)
            formulae = ['u = [u1,u2,u3,u4]', 'c = uG mod 2', 's = rHᵀ mod 2', 's → 错误位置', '取前4位恢复信息']
            for index, formula in enumerate(formulae):
                y = 2.05 + index * 0.68
                add_round_rect(slide, Inches(7.55), Inches(y), Inches(4.65), Inches(0.48), colors['white'], RGBColor(203, 213, 225))
                add_textbox(slide, Inches(7.78), Inches(y + 0.1), Inches(4.15), Inches(0.22), formula, 13, colors['dark'], index in (1, 2), PP_ALIGN.CENTER)
                if index < len(formulae) - 1:
                    add_textbox(slide, Inches(9.72), Inches(y + 0.45), Inches(0.35), Inches(0.2), '↓', 14, colors['cyan'], True, PP_ALIGN.CENTER)
        elif 'zf' in lower_title:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), 'Zero Forcing 视图', 17, colors['navy'], True, PP_ALIGN.CENTER)
            add_round_rect(slide, Inches(7.6), Inches(2.2), Inches(4.55), Inches(0.58), colors['white'], RGBColor(203, 213, 225))
            add_textbox(slide, Inches(7.8), Inches(2.36), Inches(4.15), Inches(0.2), 'channel * taps ≈ δ[n]', 16, colors['blue'], True, PP_ALIGN.CENTER)
            add_round_rect(slide, Inches(7.6), Inches(3.05), Inches(1.0), Inches(1.3), colors['cyan'])
            add_textbox(slide, Inches(7.68), Inches(3.46), Inches(0.84), Inches(0.28), 'A', 26, colors['white'], True, PP_ALIGN.CENTER)
            add_textbox(slide, Inches(8.75), Inches(3.55), Inches(0.35), Inches(0.2), '×', 22, colors['slate'], True, PP_ALIGN.CENTER)
            add_round_rect(slide, Inches(9.2), Inches(3.05), Inches(1.0), Inches(1.3), colors['green'])
            add_textbox(slide, Inches(9.28), Inches(3.46), Inches(0.84), Inches(0.28), 'w', 26, colors['white'], True, PP_ALIGN.CENTER)
            add_textbox(slide, Inches(10.35), Inches(3.55), Inches(0.35), Inches(0.2), '≈', 22, colors['slate'], True, PP_ALIGN.CENTER)
            add_round_rect(slide, Inches(10.8), Inches(3.05), Inches(1.0), Inches(1.3), colors['amber'])
            add_textbox(slide, Inches(10.88), Inches(3.46), Inches(0.84), Inches(0.28), 'd', 26, colors['white'], True, PP_ALIGN.CENTER)
        elif 'lms' in lower_title:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), 'LMS 自适应闭环', 17, colors['navy'], True, PP_ALIGN.CENTER)
            loop_items = [('x[n]', 7.65, 2.35, colors['blue']), ('y=w·x', 10.3, 2.35, colors['cyan']), ('e=d-y', 10.3, 4.0, colors['amber']), ('w←w+μex', 7.65, 4.0, colors['green'])]
            for text, x, y, fill in loop_items:
                add_round_rect(slide, Inches(x), Inches(y), Inches(1.65), Inches(0.68), fill)
                add_textbox(slide, Inches(x + 0.08), Inches(y + 0.18), Inches(1.49), Inches(0.24), text, 13, colors['white'], True, PP_ALIGN.CENTER)
            for arrow, x, y in [('→', 9.48, 2.52), ('↓', 10.95, 3.18), ('←', 9.48, 4.17), ('↑', 8.22, 3.18)]:
                add_textbox(slide, Inches(x), Inches(y), Inches(0.45), Inches(0.28), arrow, 22, colors['slate'], True, PP_ALIGN.CENTER)
        elif '环境' in title or '配置' in title:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), '环境配置流程', 17, colors['navy'], True, PP_ALIGN.CENTER)
            steps = ['git clone 仓库', 'python -m venv .venv', 'pip install -r requirements.txt', 'python src/test_environment.py', '✅ 环境验证通过']
            step_colors = [colors['blue'], colors['cyan'], colors['green'], colors['amber'], colors['slate']]
            for idx, (step, sc) in enumerate(zip(steps, step_colors)):
                sy = 2.1 + idx * 0.72
                add_round_rect(slide, Inches(7.55), Inches(sy), Inches(4.65), Inches(0.52), colors['white'], RGBColor(203, 213, 225))
                add_round_rect(slide, Inches(7.55), Inches(sy), Inches(0.12), Inches(0.52), sc, radius=False)
                add_textbox(slide, Inches(7.82), Inches(sy + 0.13), Inches(4.2), Inches(0.22), step, 12, colors['dark'], idx == 4)
                if idx < 4:
                    add_textbox(slide, Inches(9.73), Inches(sy + 0.49), Inches(0.35), Inches(0.2), '↓', 12, colors['muted'], True, PP_ALIGN.CENTER)
        elif '仓库' in title or '结构' in title or '入口' in title:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), '关键文件位置', 17, colors['navy'], True, PP_ALIGN.CENTER)
            files = [('src/', 'part1_channel_coding.py', colors['blue']), ('src/', 'part2_equalization.py', colors['green']), ('grading/', 'test_part1_coding.py', colors['cyan']), ('grading/', 'test_part2_equalization.py', colors['amber']), ('results/', 'BER曲线 / 均衡效果图', colors['slate'])]
            for idx, (folder, fname, fc) in enumerate(files):
                fy = 2.1 + idx * 0.72
                add_round_rect(slide, Inches(7.55), Inches(fy), Inches(4.65), Inches(0.52), colors['white'], RGBColor(203, 213, 225))
                add_round_rect(slide, Inches(7.55), Inches(fy), Inches(0.12), Inches(0.52), fc, radius=False)
                add_textbox(slide, Inches(7.82), Inches(fy + 0.03), Inches(1.0), Inches(0.2), folder, 9, colors['muted'])
                add_textbox(slide, Inches(7.82), Inches(fy + 0.23), Inches(4.2), Inches(0.22), fname, 12, colors['dark'], True)
        elif 'ber' in lower_title or '编码增益' in title or 'ber对比' in lower_title or '增益' in title:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), 'BER 编码增益示意', 17, colors['navy'], True, PP_ALIGN.CENTER)
            add_textbox(slide, Inches(7.55), Inches(2.05), Inches(4.65), Inches(0.28), '翻转概率 p →', 11, colors['muted'], False, PP_ALIGN.RIGHT)
            rows = [('p = 0.001', '0.0010', '0.0000', colors['green']), ('p = 0.01', '0.0100', '0.0021', colors['cyan']), ('p = 0.05', '0.0500', '0.0120', colors['amber']), ('p = 0.10', '0.1000', '0.0464', colors['red'])]
            add_round_rect(slide, Inches(7.55), Inches(2.38), Inches(4.65), Inches(0.36), colors['navy'], radius=False)
            add_textbox(slide, Inches(7.65), Inches(2.44), Inches(1.35), Inches(0.2), '翻转概率', 9, colors['white'], True, PP_ALIGN.CENTER)
            add_textbox(slide, Inches(9.15), Inches(2.44), Inches(1.45), Inches(0.2), '未编码 BER', 9, colors['white'], True, PP_ALIGN.CENTER)
            add_textbox(slide, Inches(10.68), Inches(2.44), Inches(1.45), Inches(0.2), '编码后 BER', 9, colors['white'], True, PP_ALIGN.CENTER)
            for ridx, (prob, unc, cod, rc) in enumerate(rows):
                ry = 2.82 + ridx * 0.56
                bg = colors['light'] if ridx % 2 == 0 else colors['white']
                add_round_rect(slide, Inches(7.55), Inches(ry), Inches(4.65), Inches(0.46), bg, RGBColor(203, 213, 225))
                add_textbox(slide, Inches(7.65), Inches(ry + 0.13), Inches(1.35), Inches(0.2), prob, 11, colors['dark'], False, PP_ALIGN.CENTER)
                add_textbox(slide, Inches(9.15), Inches(ry + 0.13), Inches(1.45), Inches(0.2), unc, 11, colors['red'], True, PP_ALIGN.CENTER)
                add_textbox(slide, Inches(10.68), Inches(ry + 0.13), Inches(1.45), Inches(0.2), cod, 11, rc, True, PP_ALIGN.CENTER)
        elif '多径' in title or 'isi' in lower_title or '数值例子' in title:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), 'ISI 数值计算过程', 17, colors['navy'], True, PP_ALIGN.CENTER)
            add_textbox(slide, Inches(7.55), Inches(2.05), Inches(4.65), Inches(0.28), 'h = [0.9, 0.35, −0.25]  ×  x = [+1, −1, +1]', 11, colors['slate'], True, PP_ALIGN.CENTER)
            calcs = [('y[0]', '= 0.9×(+1)', '= 0.90', colors['green']), ('y[1]', '= 0.9×(−1)+0.35×(+1)', '= −0.55', colors['amber']), ('y[2]', '= 0.9×(+1)+0.35×(−1)−0.25×(+1)', '= 0.30 ⚠', colors['red'])]
            for cidx, (yn, expr, result, rc) in enumerate(calcs):
                cy = 2.5 + cidx * 1.0
                add_round_rect(slide, Inches(7.55), Inches(cy), Inches(4.65), Inches(0.82), colors['white'], RGBColor(203, 213, 225))
                add_textbox(slide, Inches(7.72), Inches(cy + 0.06), Inches(0.62), Inches(0.28), yn, 14, colors['navy'], True)
                add_textbox(slide, Inches(8.42), Inches(cy + 0.06), Inches(3.65), Inches(0.28), expr, 11, colors['slate'])
                add_textbox(slide, Inches(8.42), Inches(cy + 0.42), Inches(3.65), Inches(0.28), result, 13, rc, True)
            add_textbox(slide, Inches(7.55), Inches(5.62), Inches(4.65), Inches(0.28), 'x[2]=+1 本应接近 +1，ISI 使其仅为 0.30', 10, colors['red'], False, PP_ALIGN.CENTER)
        elif '总结' in title or '链路' in title:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), '完整无线链路', 17, colors['navy'], True, PP_ALIGN.CENTER)
            chain = [('信息比特', colors['blue']), ('Hamming 编码', colors['cyan']), ('BPSK 调制', colors['green']), ('多径信道', colors['red']), ('ZF/LMS 均衡', colors['green']), ('BPSK 解调', colors['cyan']), ('Hamming 译码', colors['blue'])]
            for ci, (label, lc) in enumerate(chain):
                cx = 7.45 + (ci % 4) * 1.48
                cy = 2.15 + (ci // 4) * 1.45
                add_round_rect(slide, Inches(cx), Inches(cy), Inches(1.35), Inches(0.52), lc)
                add_textbox(slide, Inches(cx + 0.06), Inches(cy + 0.12), Inches(1.23), Inches(0.24), label, 9, colors['white'], True, PP_ALIGN.CENTER)
                if ci < len(chain) - 1:
                    ax = cx + 1.35 if ci % 4 < 3 else 7.45 + 1.35 * 4 + 0.5
                    ay = cy + 0.22 if ci % 4 < 3 else cy + 0.52
                    sym = '→' if ci % 4 < 3 else '↓'
                    add_textbox(slide, Inches(ax - 0.1), Inches(ay - 0.05), Inches(0.35), Inches(0.22), sym, 14, colors['muted'], True, PP_ALIGN.CENTER)
        elif '运行' in title or '提交' in title:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), '操作流程', 17, colors['navy'], True, PP_ALIGN.CENTER)
            for si, bullet in enumerate(bullets[:5]):
                sy = 2.1 + si * 0.7
                add_round_rect(slide, Inches(7.55), Inches(sy), Inches(0.38), Inches(0.38), colors['blue'])
                add_textbox(slide, Inches(7.58), Inches(sy + 0.08), Inches(0.32), Inches(0.16), str(si + 1), 10, colors['white'], True, PP_ALIGN.CENTER)
                add_textbox(slide, Inches(8.1), Inches(sy + 0.03), Inches(3.95), Inches(0.28), bullet, 12, colors['dark'])
        else:
            add_textbox(slide, Inches(7.45), Inches(1.55), Inches(4.9), Inches(0.35), '本页关键词', 17, colors['navy'], True, PP_ALIGN.CENTER)
            for ki, bullet in enumerate(bullets[:4]):
                ky = 2.2 + ki * 0.78
                add_badge(slide, 7.65, ky, f'{ki + 1}', [colors['blue'], colors['cyan'], colors['green'], colors['amber']][ki], width=0.55)
                add_textbox(slide, Inches(8.35), Inches(ky + 0.02), Inches(3.75), Inches(0.32), bullet[:34], 13, colors['dark'], ki == 0)

    def add_part_divider(slide, part_num, part_title, sub_items):
        """Render a full-bleed section-break slide for Part 1 / Part 2."""
        part_colors = {1: (colors['blue'], RGBColor(30, 58, 138)), 2: (colors['green'], RGBColor(6, 78, 59))}
        accent, bg2 = part_colors.get(part_num, (colors['cyan'], RGBColor(30, 58, 138)))
        add_round_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), colors['navy'], radius=False)
        add_round_rect(slide, Inches(0), Inches(0), Inches(5.5), Inches(7.5), bg2, radius=False)
        add_round_rect(slide, Inches(0), Inches(0), Inches(0.14), Inches(7.5), accent, radius=False)
        add_textbox(slide, Inches(0.5), Inches(1.4), Inches(4.6), Inches(0.55), f'Part {part_num}', 38, accent, True, PP_ALIGN.LEFT)
        add_round_rect(slide, Inches(0.5), Inches(2.05), Inches(4.1), Inches(0.06), accent, radius=False)
        add_textbox(slide, Inches(0.5), Inches(2.25), Inches(4.6), Inches(1.2), part_title, 24, colors['white'], True, PP_ALIGN.LEFT)
        for si, sub in enumerate(sub_items[:5]):
            sy = 3.7 + si * 0.6
            add_round_rect(slide, Inches(0.52), Inches(sy), Inches(0.28), Inches(0.28), accent)
            add_textbox(slide, Inches(0.96), Inches(sy + 0.02), Inches(4.1), Inches(0.24), sub, 12, colors['light'])
        add_textbox(slide, Inches(6.2), Inches(1.8), Inches(6.5), Inches(0.5), '无线通信技术实验 02', 20, colors['muted'], False, PP_ALIGN.LEFT)
        add_round_rect(slide, Inches(6.2), Inches(2.4), Inches(6.5), Inches(0.05), colors['muted'], radius=False)
        add_textbox(slide, Inches(6.2), Inches(2.6), Inches(6.5), Inches(0.9), '信道编码与信道均衡综合实验', 28, colors['white'], True, PP_ALIGN.LEFT)
        if part_num == 1:
            badges = [('Hamming(7,4)', colors['blue']), ('ZF 均衡', colors['cyan']), ('LMS 均衡', colors['green'])]
        else:
            badges = [('原理 + 例子', colors['green']), ('数值计算', colors['amber']), ('代码映射', colors['cyan'])]
        for bi, (bl, bc) in enumerate(badges):
            add_badge(slide, 6.2 + bi * 2.35, 5.4, bl, bc, width=2.1)

    slides = []
    current = None
    for line in Path('materials/lecture_slides_outline.md').read_text(encoding='utf-8').splitlines():
        if line.startswith('## Slide'):
            if current:
                slides.append(current)
            raw_title = line.split('：', 1)[-1] if '：' in line else line[3:]
            is_part = raw_title.startswith('[Part')
            current = {'title': raw_title, 'bullets': [], 'notes': [], 'is_part': is_part}
        elif current and line.startswith('- '):
            current['bullets'].append(line[2:])
        elif current and line.strip() and not line.startswith('#') and not line.endswith('：'):
            current['notes'].append(line.strip())
    if current:
        slides.append(current)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    total = len(slides)
    for index, item in enumerate(slides):
        layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(layout)

        # ── Cover ──────────────────────────────────────────────────────────
        if index == 0:
            cover_course = item['bullets'][0] if item['bullets'] else '无线通信技术实验'
            cover_title = item['bullets'][1] if len(item['bullets']) > 1 else item['title']
            cover_subtitle = item['bullets'][2] if len(item['bullets']) > 2 else 'Hamming(7,4) 编码 + ZF/LMS 均衡'
            cover_footer = item['bullets'][3] if len(item['bullets']) > 3 else '代码补全、仿真验证、GitHub 自动评分'
            add_round_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), colors['navy'], radius=False)
            add_round_rect(slide, Inches(0.75), Inches(0.75), Inches(11.85), Inches(5.95), RGBColor(30, 58, 138), RGBColor(30, 58, 138))
            add_textbox(slide, Inches(1.25), Inches(1.1), Inches(10.9), Inches(0.5), cover_course, 22, colors['cyan'], True, PP_ALIGN.CENTER)
            add_textbox(slide, Inches(1.25), Inches(2.05), Inches(10.9), Inches(0.82), cover_title, 36, colors['white'], True, PP_ALIGN.CENTER)
            add_textbox(slide, Inches(1.25), Inches(3.05), Inches(10.9), Inches(0.38), cover_subtitle, 20, colors['light'], True, PP_ALIGN.CENTER)
            for badge_index, badge in enumerate(item['bullets'][1:4]):
                add_badge(slide, 2.1 + badge_index * 2.35, 4.35, badge, [colors['blue'], colors['green'], colors['amber'], colors['cyan']][badge_index], width=2.05)
            add_textbox(slide, Inches(1.25), Inches(6.25), Inches(10.9), Inches(0.28), f'两部分讲解：实验流程 + 技术原理 · {cover_footer}', 14, colors['light'], False, PP_ALIGN.CENTER)
            continue

        # ── Part divider ───────────────────────────────────────────────────
        if item.get('is_part'):
            import re
            m = re.match(r'\[Part(\d+)\]\s*(.*)', item['title'])
            part_num = int(m.group(1)) if m else 1
            part_title_text = m.group(2) if m else item['title']
            add_part_divider(slide, part_num, part_title_text, item['bullets'])
            continue

        # ── Normal content slide ───────────────────────────────────────────
        add_round_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), colors['white'], radius=False)
        add_round_rect(slide, Inches(0), Inches(0.72), Inches(13.333), Inches(0.06), colors['cyan'], radius=False)
        add_header(slide, item['title'], index + 1, total)
        add_bullet_cards(slide, item['bullets'])
        add_visual_panel(slide, item['title'], item['bullets'], index + 1)
        if item['notes']:
            note_text = item['notes'][0]
            add_round_rect(slide, Inches(0.65), Inches(6.55), Inches(12.0), Inches(0.38), RGBColor(239, 246, 255), RGBColor(191, 219, 254))
            add_textbox(slide, Inches(0.82), Inches(6.64), Inches(11.65), Inches(0.18), f'讲解提示：{note_text[:86]}', 9, colors['slate'])
        add_footer(slide)
    presentation.save('materials/lecture_slides.pptx')


def main():
    Path('materials').mkdir(exist_ok=True)
    generate_word()
    generate_ppt()
    print('已生成 materials/teacher_lab_guide.docx 和 materials/lecture_slides.pptx')


if __name__ == '__main__':
    main()
